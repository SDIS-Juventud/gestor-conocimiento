# -*- coding: utf-8 -*-
# Genera la presentacion de buenas practicas del GECO chat en NotebookLM,
# como material de trabajo para revisar con David. Estilo institucional sobrio
# (sin gradientes, sin sombras dramaticas, paleta del gestor).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(BASE, "imagenes")
SALIDA = os.path.join(BASE, "ppt", "GECO_chat_buenas_practicas_NotebookLM.pptx")

# Paleta del gestor
SLATE  = RGBColor(0x2F, 0x3E, 0x3C)   # verde-gris oscuro (titulos)
ACCENT = RGBColor(0x66, 0x3A, 0x93)   # morado (acentos)
INK    = RGBColor(0x33, 0x33, 0x33)   # texto cuerpo
GRAY   = RGBColor(0x8a, 0x8a, 0x8a)   # pie de pagina
CREAM  = RGBColor(0xF8, 0xF4, 0xE1)   # fondo portada
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Figtree"        # cuerpo, igual que el gestor
TITLE_FONT = "Anton"    # titulos, igual que el gestor

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def _set_font(run, size, color, bold=False, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rule(slide, x, y, w, color=ACCENT, h=Pt(2.4)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def footer(slide, page):
    tb, tf = textbox(slide, Inches(0.55), Inches(7.02), Inches(9), Inches(0.32))
    r = tf.paragraphs[0].add_run()
    r.text = "GECO chat · NotebookLM   |   Subdirección para la Juventud · SDIS"
    _set_font(r, 9, GRAY)
    tb2, tf2 = textbox(slide, Inches(11.6), Inches(7.02), Inches(1.2), Inches(0.32))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(page)
    _set_font(r2, 9, GRAY)

def content_slide(kicker, title, bullets, page):
    """bullets: lista de dict {lead, text, level}. lead va en negrita."""
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide, WHITE)
    # Kicker
    tb, tf = textbox(slide, Inches(0.55), Inches(0.5), Inches(11.5), Inches(0.35))
    r = tf.paragraphs[0].add_run(); r.text = kicker.upper()
    _set_font(r, 11, ACCENT, bold=True)
    # Titulo
    tb, tf = textbox(slide, Inches(0.55), Inches(0.85), Inches(12.2), Inches(1.0))
    r = tf.paragraphs[0].add_run(); r.text = title
    _set_font(r, 26, SLATE, name=TITLE_FONT)
    # Regla
    rule(slide, Inches(0.57), Inches(1.78), Inches(1.5))
    # Cuerpo
    tb, tf = textbox(slide, Inches(0.6), Inches(2.05), Inches(12.1), Inches(4.7))
    first = True
    for b in bullets:
        lvl = b.get("level", 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(9 if lvl == 0 else 4)
        p.line_spacing = 1.08
        bullet = "—  " if lvl == 0 else "·  "
        rb = p.add_run(); rb.text = bullet
        _set_font(rb, 15 if lvl == 0 else 13, ACCENT if lvl == 0 else GRAY, bold=True)
        if b.get("lead"):
            rl = p.add_run(); rl.text = b["lead"] + "  "
            _set_font(rl, 15 if lvl == 0 else 13, SLATE, bold=True)
        if b.get("text"):
            rt = p.add_run(); rt.text = b["text"]
            _set_font(rt, 15 if lvl == 0 else 13, INK)
        pPr = p._p.get_or_add_pPr()
        marL = Inches(0.0) if lvl == 0 else Inches(0.45)
        pPr.set('marL', str(int(marL)))
    footer(slide, page)
    return slide

# ---------------- Portada ----------------
s = prs.slides.add_slide(BLANK)
fill_bg(s, CREAM)
rule(s, Inches(0.9), Inches(2.55), Inches(2.2), ACCENT, h=Pt(3.2))
tb, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "SUBDIRECCIÓN PARA LA JUVENTUD · SDIS"
_set_font(r, 13, ACCENT, bold=True)
tb, tf = textbox(s, Inches(0.88), Inches(2.75), Inches(11.6), Inches(1.8))
r = tf.paragraphs[0].add_run(); r.text = "GECO chat: buenas prácticas en NotebookLM"
_set_font(r, 40, SLATE, name=TITLE_FONT)
tb, tf = textbox(s, Inches(0.9), Inches(4.3), Inches(10.8), Inches(0.9))
r = tf.paragraphs[0].add_run()
r.text = "Cómo organizar la documentación para que el asistente responda bien"
_set_font(r, 19, INK)
tb, tf = textbox(s, Inches(0.9), Inches(6.4), Inches(8.5), Inches(0.4))
r = tf.paragraphs[0].add_run()
r.text = "Material de trabajo  ·  Junio de 2026  ·  Equipo de Analítica e Innovación"
_set_font(r, 12, GRAY)
try:
    s.shapes.add_picture(os.path.join(IMG, "Footer1.png"), Inches(10.0), Inches(6.3), height=Inches(0.5))
    s.shapes.add_picture(os.path.join(IMG, "Footer2.png"), Inches(11.5), Inches(6.3), height=Inches(0.5))
except Exception as e:
    print("  (logos no agregados:", e, ")")

# ---------------- Contenido ----------------
content_slide(
    "Punto de partida", "Dos componentes que se complementan",
    [
        {"lead":"El gestor (sitio en GitHub).", "text":"Páginas por servicio con su contexto, protocolos y enlaces a documentos."},
        {"lead":"El GECO chat (NotebookLM).", "text":"Asistente que responde preguntas en lenguaje natural sobre esos mismos documentos, citando siempre la fuente."},
        {"lead":"Idea clave.", "text":"El chat solo sabe lo que se le sube. Su calidad depende de cómo organizamos y cargamos la documentación."},
    ], 2)

content_slide(
    "La pregunta", "¿Meter más documentos de un servicio que de otro daña el chat?",
    [
        {"lead":"Respuesta corta: no.", "text":""},
        {"lead":"No reparte por cuotas.", "text":"NotebookLM funciona por relevancia (RAG): ante una pregunta, busca los fragmentos más parecidos entre todas las fuentes y responde citándolos."},
        {"lead":"No se opaca un servicio.", "text":"Uno con 3 documentos no queda tapado por otro con 20, siempre que la pregunta sea específica de ese servicio."},
    ], 3)

content_slide(
    "Lo que sí importa", "No es la cantidad, son estas cuatro cosas",
    [
        {"lead":"Cobertura, no cantidad.", "text":"El chat solo responde lo que existe. Tres documentos buenos valen más que veinte repetitivos."},
        {"lead":"Preguntas específicas.", "text":"Indicar el servicio evita respuestas mezcladas cuando el tema aplica a varios."},
        {"lead":"Seleccionar fuentes.", "text":"En NotebookLM se pueden marcar solo las fuentes de un servicio antes de preguntar. Es el mejor seguro contra el desbalance."},
        {"lead":"Límite de fuentes.", "text":"El plan gratuito permite 50 fuentes por notebook. Cuidar el cupo, no inflar con relleno."},
    ], 4)

content_slide(
    "Matiz técnico", "El chat lee el contenido, no el nombre del archivo",
    [
        {"lead":"Lee el texto de adentro.", "text":"NotebookLM no usa el prefijo del nombre para enrutar; usa lo que dice el documento."},
        {"lead":"Los códigos sirven a las personas.", "text":"03, 06, 07… ayudan a organizar y a filtrar fuentes, no a que el chat decida."},
        {"lead":"Consecuencia.", "text":"Importa mucho cómo está escrito y estructurado cada documento."},
    ], 5)

content_slide(
    "¿Índice más detallado?", "Sí, pero distinguiendo dos propósitos",
    [
        {"lead":"Para las personas (trazabilidad).", "text":"El índice maestro de nomenclatura, mantenido al día. Aquí más detalle siempre ayuda."},
        {"lead":"Para el chat (mejores respuestas).", "text":"Una ficha-encabezado corta al inicio de cada documento: qué es, servicio, fecha, versión y para qué sirve."},
        {"text":"Ese texto se indexa, da contexto al fragmento y hace las citas más precisas.", "level":1},
        {"text":"Es más útil que un índice externo larguísimo que el chat ni siquiera lee.", "level":1},
    ], 6)

content_slide(
    "Otras prácticas útiles (1/2)", "Calidad y formato de las fuentes",
    [
        {"lead":"Texto seleccionable, no imágenes.", "text":"NotebookLM no lee PDF escaneados sin OCR. Subir documentos con texto real."},
        {"lead":"Sin duplicados ni versiones viejas.", "text":"Si conviven v1 y v2, el chat puede citar la vieja. Reemplazar, no acumular."},
        {"lead":"Datos en texto.", "text":"Como Power BI no se enlaza, los resúmenes deben traer las cifras clave escritas, no solo remitir al tablero."},
        {"lead":"Nombres sin tildes ni espacios.", "text":"Usar guion bajo (03_SDIS_JUV_7940) para que no se rompan al subir."},
        {"lead":"Documentos autocontenidos.", "text":"Que cada uno se entienda solo, sin depender de otro."},
    ], 7)

content_slide(
    "Otras prácticas útiles (2/2)", "Comportamiento y gobernanza",
    [
        {"lead":"Reglas claras del asistente.", "text":"En 01_reglas_generales: citar siempre la fuente, no inventar, decir cuándo no hay información."},
        {"lead":"Solo fuentes oficiales y validadas.", "text":"El chat trata todo lo que se sube como verdad. No meter borradores sin revisar."},
        {"lead":"Un responsable de cargue.", "text":"Una persona o rol que suba y revise, siguiendo la nomenclatura."},
        {"lead":"Rutina de actualización.", "text":"Cada documento nuevo o cambio entra con la nomenclatura y reemplaza la versión anterior."},
        {"lead":"Probarlo con preguntas reales.", "text":"Armar un set de preguntas del día a día y ajustar según los fallos."},
    ], 8)

content_slide(
    "Nomenclatura actual", "Recordatorio de los códigos vigentes",
    [
        {"lead":"01–02", "text":"Introducción al GECO chat (reglas generales, qué es el gestor)."},
        {"lead":"03", "text":"Generalidades SDIS y Subdirección (Decreto 647 de 2025, proyecto 7940)."},
        {"lead":"04", "text":"Páginas del gestor (cada HTML, por enlace directo)."},
        {"lead":"05", "text":"Datos: SIRBE (Ficha SIRBE + resúmenes escritos de los tableros Power BI)."},
        {"lead":"06", "text":"Jóvenes con Oportunidades (pendiente: portafolio SDIS)."},
        {"lead":"07 · 08 · 09", "text":"Casas de Juventud · Parche Seguro · Servicio Forjar Restaurativo."},
    ], 9)

content_slide(
    "Para revisar con David", "Decisiones y próximos pasos",
    [
        {"lead":"Aprobar la ficha-encabezado estándar.", "text":"Una plantilla única para pegar al inicio de cada documento."},
        {"lead":"Completar cobertura faltante.", "text":"Por ejemplo, el portafolio SDIS de Jóvenes con Oportunidades."},
        {"lead":"Definir responsable y rutina de cargue.", "text":"Quién sube, quién revisa y cada cuánto."},
        {"lead":"Vigilar el límite de 50 fuentes.", "text":"Evaluar plan Pro si la documentación crece."},
        {"lead":"Acordar un set de preguntas de prueba.", "text":"Antes de abrir el chat al equipo."},
    ], 10)

os.makedirs(os.path.dirname(SALIDA), exist_ok=True)
prs.save(SALIDA)
print("OK ->", SALIDA)
print("Diapositivas:", len(prs.slides._sldIdLst))
